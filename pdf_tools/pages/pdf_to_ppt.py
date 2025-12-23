import html
from pathlib import Path
from tempfile import TemporaryDirectory

from PySide6.QtCore import Qt, QThread, Signal, QUrl
from PySide6.QtGui import QDesktopServices, QDragEnterEvent, QDropEvent
from PySide6.QtWidgets import QFileDialog, QHBoxLayout, QLabel, QLineEdit, QMessageBox, QProgressBar, QPushButton, QSpinBox, QVBoxLayout, QWidget


class DropArea(QLabel):
    file_selected = Signal(Path)

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.setText("点击或拖拽PDF文件到此处")
        self.setStyleSheet(
            """
            QLabel {
                border: 2px dashed #aaa;
                border-radius: 10px;
                padding: 14px;
                background-color: #f9f9f9;
                color: #555;
                font-size: 14px;
            }
            QLabel:hover {
                background-color: #eef;
                border-color: #88d;
            }
            """
        )
        self.setAcceptDrops(True)
        self.setMinimumHeight(60)
        self.setCursor(Qt.CursorShape.PointingHandCursor)

    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            for url in event.mimeData().urls():
                if url.toLocalFile().lower().endswith(".pdf"):
                    event.acceptProposedAction()
                    return
        event.ignore()

    def dropEvent(self, event: QDropEvent):
        for url in event.mimeData().urls():
            f = url.toLocalFile()
            if f.lower().endswith(".pdf"):
                self.update_file(Path(f))
                break
        event.acceptProposedAction()

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            f, _ = QFileDialog.getOpenFileName(self, "选择PDF", "", "PDF (*.pdf)")
            if f:
                self.update_file(Path(f))

    def update_file(self, path: Path):
        name = html.escape(path.name)
        full = html.escape(str(path.absolute()))
        self.setText(f"<div style='line-height:1.7'>{name}<br/>{full}</div>")
        self.setTextFormat(Qt.TextFormat.RichText)
        self.setStyleSheet(
            """
            QLabel {
                border: 2px solid #4caf50;
                border-radius: 10px;
                padding: 14px;
                background-color: #e8f5e9;
                color: #2e7d32;
                font-size: 14px;
            }
            """
        )
        self.file_selected.emit(path)


class PdfToPptThread(QThread):
    finished_signal = Signal(str)
    error_signal = Signal(str)

    def __init__(self, pdf_path: str, out_path: str, dpi: int):
        super().__init__()
        self.pdf_path = pdf_path
        self.out_path = out_path
        self.dpi = dpi

    def run(self):
        try:
            import fitz
            from pptx import Presentation

            out = Path(self.out_path)
            out.parent.mkdir(parents=True, exist_ok=True)

            doc = fitz.open(self.pdf_path)
            zoom = self.dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)

            prs = Presentation()
            blank_layout = prs.slide_layouts[6]

            slide_w = prs.slide_width
            slide_h = prs.slide_height

            with TemporaryDirectory() as td:
                td_path = Path(td)
                for i in range(doc.page_count):
                    page = doc.load_page(i)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    png = td_path / f"page_{i+1:03d}.png"
                    pix.save(str(png))

                    slide = prs.slides.add_slide(blank_layout)
                    pic = slide.shapes.add_picture(str(png), 0, 0)

                    img_w = pic.width
                    img_h = pic.height
                    scale = min(slide_w / img_w, slide_h / img_h)
                    new_w = int(img_w * scale)
                    new_h = int(img_h * scale)
                    left = int((slide_w - new_w) / 2)
                    top = int((slide_h - new_h) / 2)
                    pic.left = left
                    pic.top = top
                    pic.width = new_w
                    pic.height = new_h

            prs.save(str(out))
            self.finished_signal.emit(str(out))
        except Exception as e:
            self.error_signal.emit(str(e))


class PdfToPptPage(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.current_pdf_path: Path | None = None
        self.thread: PdfToPptThread | None = None

        lay = QVBoxLayout(self)
        lay.setContentsMargins(12, 12, 12, 12)
        lay.setSpacing(12)

        title = QLabel("PDF转PPT")
        title.setAlignment(Qt.AlignmentFlag.AlignLeft)
        lay.addWidget(title)

        self.drop_area = DropArea()
        self.drop_area.file_selected.connect(self.on_file_selected)
        lay.addWidget(self.drop_area)

        r1 = QHBoxLayout()
        self.dpi_spin = QSpinBox()
        self.dpi_spin.setMinimum(72)
        self.dpi_spin.setMaximum(600)
        self.dpi_spin.setValue(150)
        r1.addWidget(QLabel("DPI"))
        r1.addWidget(self.dpi_spin)
        r1.addStretch(1)
        lay.addLayout(r1)

        r2 = QHBoxLayout()
        self.output_edit = QLineEdit()
        self.output_edit.setReadOnly(True)
        self.output_edit.setPlaceholderText("输出PPTX路径")
        self.output_btn = QPushButton("选择输出文件")
        r2.addWidget(QLabel("输出"))
        r2.addWidget(self.output_edit, 1)
        r2.addWidget(self.output_btn)
        lay.addLayout(r2)

        r3 = QHBoxLayout()
        self.open_btn = QPushButton("打开文件")
        self.open_btn.setEnabled(False)
        r3.addStretch(1)
        r3.addWidget(self.open_btn)
        lay.addLayout(r3)

        self.progress = QProgressBar()
        self.progress.setRange(0, 0)
        self.progress.hide()
        lay.addWidget(self.progress)

        self.convert_btn = QPushButton("开始转换")
        lay.addWidget(self.convert_btn, 0, Qt.AlignmentFlag.AlignLeft)

        self.output_btn.clicked.connect(self.on_choose_output)
        self.open_btn.clicked.connect(self.on_open)
        self.convert_btn.clicked.connect(self.on_convert)

    def on_file_selected(self, path: Path):
        self.current_pdf_path = path
        if not self.output_edit.text().strip():
            self.output_edit.setText(str(path.with_suffix(".pptx")))

    def on_choose_output(self):
        suggested = self.output_edit.text().strip()
        f, _ = QFileDialog.getSaveFileName(self, "选择输出PPTX", suggested or "", "PowerPoint (*.pptx)")
        if f:
            if not f.lower().endswith(".pptx"):
                f += ".pptx"
            self.output_edit.setText(f)
            self.open_btn.setEnabled(Path(f).exists())

    def on_open(self):
        p = self.output_edit.text().strip()
        if p and Path(p).exists():
            QDesktopServices.openUrl(QUrl.fromLocalFile(p))
        else:
            QMessageBox.warning(self, "提示", "文件不存在")

    def on_convert(self):
        inp = self.current_pdf_path
        if not inp or not inp.is_file():
            QMessageBox.warning(self, "提示", "请选择输入PDF")
            return
        out = self.output_edit.text().strip()
        if not out:
            QMessageBox.warning(self, "提示", "请选择输出文件")
            return
        dpi = int(self.dpi_spin.value())

        self.progress.show()
        self.convert_btn.setEnabled(False)
        self.open_btn.setEnabled(False)
        self.thread = PdfToPptThread(str(inp), out, dpi)
        self.thread.finished_signal.connect(self.on_finished)
        self.thread.error_signal.connect(self.on_error)
        self.thread.finished.connect(self.on_thread_done)
        self.thread.start()

    def on_finished(self, out: str):
        self.output_edit.setText(out)
        self.open_btn.setEnabled(Path(out).exists())
        QMessageBox.information(self, "完成", out)

    def on_error(self, msg: str):
        QMessageBox.critical(self, "错误", msg)

    def on_thread_done(self):
        self.progress.hide()
        self.convert_btn.setEnabled(True)
